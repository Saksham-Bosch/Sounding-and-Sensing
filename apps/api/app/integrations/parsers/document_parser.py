from pathlib import Path

import docx
import fitz  # PyMuPDF
import httpx
import pandas as pd
import pptx
from bs4 import BeautifulSoup

from app.integrations.ocr.client import extract_text_with_ocr


async def parse_document(file_path: Path, mime_type: str) -> str:
    """Extracts text deterministically, falling back to OCR for scanned PDFs."""
    extracted_text = ""
    suffix = file_path.suffix.lower()

    try:
        # 1. Images -> Direct to OCR
        if suffix in [".png", ".jpg", ".jpeg"]:
            return await extract_text_with_ocr(file_path)

        # 2. PDF Handling (Digital vs Scanned)
        elif suffix == ".pdf":
            doc = fitz.open(file_path)
            for page in doc:
                extracted_text += page.get_text() + "\n"

            # If PyMuPDF found almost no text, assume it's a scanned PDF
            if len(extracted_text.strip()) < 50 and len(doc) > 0:
                print("PDF appears scanned. Running OCR on first page...")
                first_page = doc[0]
                pix = first_page.get_pixmap()
                temp_img_path = file_path.with_suffix(".temp.png")
                pix.save(temp_img_path)
                extracted_text = await extract_text_with_ocr(temp_img_path)
                temp_img_path.unlink(missing_ok=True)  # Clean up temp image
            doc.close()

        # 3. Word Documents
        elif suffix == ".docx":
            doc = docx.Document(file_path)
            extracted_text = "\n".join([para.text for para in doc.paragraphs])

        # 4. PowerPoint Presentations
        elif suffix == ".pptx":
            prs = pptx.Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        extracted_text += shape.text + "\n"

        # 5. Excel Spreadsheets
        elif suffix in [".xlsx", ".csv"]:
            df = pd.read_excel(file_path) if suffix == ".xlsx" else pd.read_csv(file_path)
            extracted_text = df.to_string(index=False)

        # 6. Plain Text
        elif suffix == ".txt":
            with open(file_path, "r", encoding="utf-8") as f:
                extracted_text = f.read()
        else:
            extracted_text = f"[Unsupported document type: {suffix}]"

    except Exception as e:
        extracted_text = f"[Error extracting text: {str(e)}]"

    return extracted_text.strip()


async def parse_url(url: str) -> str:
    """Scrapes clean text from a webpage."""
    async with httpx.AsyncClient(timeout=15.0, verify=False) as client:
        try:
            response = await client.get(url)
            response.raise_for_status()
            soup = BeautifulSoup(response.text, "html.parser")
            # Remove scripts and styles
            for script in soup(["script", "style"]):
                script.extract()
            text = soup.get_text(separator="\n")
            # Clean up blank lines
            lines = (line.strip() for line in text.splitlines())
            return "\n".join(chunk for chunk in lines if chunk)
        except Exception as e:
            return f"[Failed to scrape URL: {str(e)}]"
